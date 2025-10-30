#!/usr/bin/env python3
"""
Create a sample PowerPoint presentation for testing the application.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_sample_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Sample Presentation"
    
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.size = Pt(48)
    title_run.font.bold = True
    
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(8), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Test PowerPoint for Processing"
    
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_para.runs[0]
    subtitle_run.font.size = Pt(24)
    
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    heading_box = slide2.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(0.5)
    )
    heading_frame = heading_box.text_frame
    heading_frame.text = "Key Features"
    heading_para = heading_frame.paragraphs[0]
    heading_run = heading_para.runs[0]
    heading_run.font.size = Pt(36)
    heading_run.font.bold = True
    
    content_box = slide2.shapes.add_textbox(
        Inches(1.5), Inches(2), Inches(7), Inches(4)
    )
    content_frame = content_box.text_frame
    
    features = [
        "Automatic branding application",
        "Font standardization",
        "Color scheme consistency",
        "Layout preservation",
        "Fast processing"
    ]
    
    for feature in features:
        p = content_frame.add_paragraph()
        p.text = f"• {feature}"
        p.level = 0
        p.font.size = Pt(20)
    
    slide3 = prs.slides.add_slide(blank_slide_layout)
    
    heading3_box = slide3.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(0.5)
    )
    heading3_frame = heading3_box.text_frame
    heading3_frame.text = "Benefits"
    h3_para = heading3_frame.paragraphs[0]
    h3_run = h3_para.runs[0]
    h3_run.font.size = Pt(36)
    h3_run.font.bold = True
    
    content3_box = slide3.shapes.add_textbox(
        Inches(1.5), Inches(2), Inches(7), Inches(4)
    )
    content3_frame = content3_box.text_frame
    content3_frame.text = "Save time on formatting and focus on content. Ensure consistent branding across all presentations. Professional appearance with minimal effort."
    
    for para in content3_frame.paragraphs:
        para.font.size = Pt(20)
    
    output_path = os.path.join(os.path.dirname(__file__), 'sample_presentation.pptx')
    prs.save(output_path)
    print(f"Sample presentation created: {output_path}")
    return output_path

if __name__ == '__main__':
    create_sample_presentation()
