from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, Any
import logging

logger = logging.getLogger(__name__)


class BrandingEngine:
    BRAND_CONFIG = {
        'fonts': {
            'primary': 'Frutiger LT',
            'fallback': ['Calibri', 'Arial', 'Helvetica'],
            'title_size': 44,
            'heading_size': 32,
            'body_size': 18,
            'small_size': 14
        },
        'colors': {
            'primary': RGBColor(0, 102, 204),
            'secondary': RGBColor(102, 102, 102),
            'accent': RGBColor(255, 107, 53),
            'text_dark': RGBColor(51, 51, 51),
            'text_light': RGBColor(255, 255, 255),
            'background': RGBColor(255, 255, 255)
        },
        'spacing': {
            'title_margin_top': Inches(0.5),
            'heading_margin_top': Inches(0.3),
            'body_margin_top': Inches(0.2),
            'line_spacing': 1.2
        },
        'alignment': {
            'title': PP_ALIGN.CENTER,
            'heading': PP_ALIGN.LEFT,
            'body': PP_ALIGN.LEFT
        },
        'transparency': {
            'image_overlay': 0.85,
            'background': 1.0
        }
    }
    
    def __init__(self, custom_config: Dict[str, Any] = None):
        self.config = self.BRAND_CONFIG.copy()
        if custom_config:
            self.config.update(custom_config)
    
    def apply_branding_to_shape(self, shape, shape_type: str = 'body'):
        if not shape.has_text_frame:
            return
        
        try:
            text_frame = shape.text_frame
            
            for paragraph in text_frame.paragraphs:
                self.apply_paragraph_formatting(paragraph, shape_type)
                
                for run in paragraph.runs:
                    self.apply_run_formatting(run, shape_type)
            
            logger.debug(f"Applied branding to shape: {shape.name}")
        except Exception as e:
            logger.error(f"Error applying branding to shape: {str(e)}")
    
    def apply_paragraph_formatting(self, paragraph, shape_type: str):
        try:
            alignment = self.config['alignment'].get(shape_type, PP_ALIGN.LEFT)
            paragraph.alignment = alignment
            paragraph.line_spacing = self.config['spacing']['line_spacing']
        except Exception as e:
            logger.warning(f"Could not apply paragraph formatting: {str(e)}")
    
    def apply_run_formatting(self, run, shape_type: str):
        try:
            font = run.font
            
            font_name = self.get_brand_font()
            font.name = font_name
            
            font_size = self.get_font_size(shape_type)
            font.size = font_size
            
            if shape_type == 'title':
                font.bold = True
                font.color.rgb = self.config['colors']['primary']
            elif shape_type == 'heading':
                font.bold = True
                font.color.rgb = self.config['colors']['text_dark']
            else:
                font.bold = False
                font.color.rgb = self.config['colors']['text_dark']
        
        except Exception as e:
            logger.warning(f"Could not apply run formatting: {str(e)}")
    
    def get_brand_font(self) -> str:
        primary = self.config['fonts']['primary']
        return primary
    
    def get_font_size(self, shape_type: str):
        size_map = {
            'title': self.config['fonts']['title_size'],
            'heading': self.config['fonts']['heading_size'],
            'body': self.config['fonts']['body_size'],
            'small': self.config['fonts']['small_size']
        }
        return Pt(size_map.get(shape_type, self.config['fonts']['body_size']))
    
    def apply_color_scheme(self, presentation):
        try:
            for slide in presentation.slides:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = self.config['colors']['background']
        except Exception as e:
            logger.warning(f"Could not apply color scheme: {str(e)}")
    
    def standardize_visual_elements(self, shape):
        try:
            if hasattr(shape, 'fill'):
                if shape.fill.type == 1:
                    pass
        except Exception as e:
            logger.debug(f"Could not standardize visual elements: {str(e)}")
    
    def apply_logo_placement(self, slide, logo_path: str = None):
        if not logo_path:
            return
        
        try:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(1.5)
            slide.shapes.add_picture(logo_path, left, top, width=width)
            logger.debug("Logo added to slide")
        except Exception as e:
            logger.error(f"Could not add logo: {str(e)}")
    
    def determine_shape_type(self, shape, slide_index: int) -> str:
        if not shape.has_text_frame:
            return 'none'
        
        text = shape.text.strip()
        
        if not text:
            return 'none'
        
        if slide_index == 0 and shape.top < Inches(2):
            return 'title'
        
        if shape.top < Inches(1.5):
            return 'title'
        
        text_length = len(text)
        if text_length < 50 and shape.width > Inches(5):
            return 'heading'
        
        font_size = None
        if shape.text_frame.paragraphs:
            first_run = shape.text_frame.paragraphs[0].runs
            if first_run:
                font_size = first_run[0].font.size
        
        if font_size and font_size > Pt(28):
            return 'title'
        elif font_size and font_size > Pt(20):
            return 'heading'
        
        return 'body'
