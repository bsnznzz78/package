from pptx import Presentation
from typing import Dict, Any, Optional
import logging
import os

logger = logging.getLogger(__name__)


class PPTParser:
    ALLOWED_EXTENSIONS = {'.pptx'}
    MAX_FILE_SIZE = 50 * 1024 * 1024
    
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.presentation = None
        self._validate_file()
    
    def _validate_file(self):
        if not os.path.exists(self.file_path):
            raise FileNotFoundError(f"File not found: {self.file_path}")
        
        file_ext = os.path.splitext(self.file_path)[1].lower()
        if file_ext not in self.ALLOWED_EXTENSIONS:
            raise ValueError(f"Invalid file format. Expected .pptx, got {file_ext}")
        
        file_size = os.path.getsize(self.file_path)
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"File too large. Max size: {self.MAX_FILE_SIZE / (1024*1024)}MB")
        
        logger.info(f"File validation passed: {self.file_path}")
    
    def load_presentation(self) -> Presentation:
        try:
            self.presentation = Presentation(self.file_path)
            logger.info(f"Loaded presentation with {len(self.presentation.slides)} slides")
            return self.presentation
        except Exception as e:
            logger.error(f"Error loading presentation: {str(e)}")
            raise ValueError(f"Could not load PowerPoint file: {str(e)}")
    
    def get_presentation_info(self) -> Dict[str, Any]:
        if not self.presentation:
            self.load_presentation()
        
        return {
            'slide_count': len(self.presentation.slides),
            'slide_width': self.presentation.slide_width,
            'slide_height': self.presentation.slide_height,
            'layouts_count': len(self.presentation.slide_layouts),
            'has_notes': any(slide.has_notes_slide for slide in self.presentation.slides)
        }
    
    def validate_presentation_structure(self) -> Dict[str, Any]:
        if not self.presentation:
            self.load_presentation()
        
        validation = {
            'is_valid': True,
            'errors': [],
            'warnings': []
        }
        
        if len(self.presentation.slides) == 0:
            validation['is_valid'] = False
            validation['errors'].append("Presentation has no slides")
        
        if len(self.presentation.slides) > 200:
            validation['warnings'].append("Presentation has more than 200 slides. Processing may take longer.")
        
        has_content = False
        for slide in self.presentation.slides:
            if slide.shapes:
                has_content = True
                break
        
        if not has_content:
            validation['warnings'].append("No content found in presentation")
        
        logger.info(f"Validation result: {validation}")
        return validation
    
    @staticmethod
    def is_valid_pptx_file(file_path: str) -> bool:
        try:
            if not os.path.exists(file_path):
                return False
            
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext not in PPTParser.ALLOWED_EXTENSIONS:
                return False
            
            Presentation(file_path)
            return True
        except Exception as e:
            logger.error(f"File validation failed: {str(e)}")
            return False
