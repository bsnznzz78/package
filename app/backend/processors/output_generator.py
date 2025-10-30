from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, Any, Optional
import logging
import os
from datetime import datetime

from .branding_engine import BrandingEngine

logger = logging.getLogger(__name__)


class OutputGenerator:
    def __init__(self, input_ppt_path: str, output_dir: str):
        self.input_ppt_path = input_ppt_path
        self.output_dir = output_dir
        self.branding_engine = BrandingEngine()
        self.input_presentation = None
        self.output_presentation = None
        
    def generate_enhanced_presentation(self, branding_config: Dict[str, Any] = None) -> str:
        try:
            if branding_config:
                self.branding_engine = BrandingEngine(branding_config)
            
            self.input_presentation = Presentation(self.input_ppt_path)
            self.output_presentation = Presentation(self.input_ppt_path)
            
            logger.info(f"Processing {len(self.output_presentation.slides)} slides")
            
            for idx, slide in enumerate(self.output_presentation.slides):
                self.process_slide(slide, idx)
            
            self.branding_engine.apply_color_scheme(self.output_presentation)
            
            output_path = self.save_presentation()
            logger.info(f"Enhanced presentation saved to: {output_path}")
            
            return output_path
        
        except Exception as e:
            logger.error(f"Error generating enhanced presentation: {str(e)}")
            raise
    
    def process_slide(self, slide, slide_index: int):
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_type = self.branding_engine.determine_shape_type(shape, slide_index)
                    self.branding_engine.apply_branding_to_shape(shape, shape_type)
                
                self.branding_engine.standardize_visual_elements(shape)
            
            logger.debug(f"Processed slide {slide_index + 1}")
        
        except Exception as e:
            logger.error(f"Error processing slide {slide_index}: {str(e)}")
            raise
    
    def save_presentation(self) -> str:
        os.makedirs(self.output_dir, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"enhanced_presentation_{timestamp}.pptx"
        output_path = os.path.join(self.output_dir, filename)
        
        self.output_presentation.save(output_path)
        
        return output_path
    
    def add_watermark(self, text: str = "Enhanced by PPT Processor"):
        try:
            for slide in self.output_presentation.slides:
                left = Inches(0.1)
                top = self.output_presentation.slide_height - Inches(0.5)
                width = Inches(3)
                height = Inches(0.3)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text
                
                paragraph = text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.size = Pt(8)
                run.font.color.rgb = self.branding_engine.config['colors']['secondary']
        
        except Exception as e:
            logger.warning(f"Could not add watermark: {str(e)}")
    
    def preserve_media_elements(self, source_slide, target_slide):
        try:
            for shape in source_slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pass
                elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    pass
        except Exception as e:
            logger.warning(f"Could not preserve media elements: {str(e)}")
    
    def validate_output(self, output_path: str) -> bool:
        try:
            if not os.path.exists(output_path):
                logger.error("Output file does not exist")
                return False
            
            test_prs = Presentation(output_path)
            
            if len(test_prs.slides) != len(self.input_presentation.slides):
                logger.error("Slide count mismatch")
                return False
            
            logger.info("Output validation passed")
            return True
        
        except Exception as e:
            logger.error(f"Output validation failed: {str(e)}")
            return False
