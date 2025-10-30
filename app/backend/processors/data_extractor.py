from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Any
import logging

logger = logging.getLogger(__name__)


class DataExtractor:
    def __init__(self, ppt_path: str):
        self.ppt_path = ppt_path
        self.presentation = Presentation(ppt_path)
        
    def extract_all_data(self) -> Dict[str, Any]:
        try:
            data = {
                'slides': self.extract_slides_data(),
                'branding': self.extract_branding_info(),
                'metadata': self.extract_metadata()
            }
            logger.info(f"Successfully extracted data from {len(data['slides'])} slides")
            return data
        except Exception as e:
            logger.error(f"Error extracting data: {str(e)}")
            raise
    
    def extract_slides_data(self) -> List[Dict[str, Any]]:
        slides_data = []
        
        for idx, slide in enumerate(self.presentation.slides):
            slide_data = {
                'index': idx,
                'layout': slide.slide_layout.name,
                'shapes': self.extract_shapes_from_slide(slide),
                'notes': self.extract_notes(slide)
            }
            slides_data.append(slide_data)
        
        return slides_data
    
    def extract_shapes_from_slide(self, slide) -> List[Dict[str, Any]]:
        shapes_data = []
        
        for shape in slide.shapes:
            shape_info = {
                'type': str(shape.shape_type),
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            
            if shape.has_text_frame:
                shape_info['text_content'] = self.extract_text_from_shape(shape)
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info['image'] = self.extract_image_info(shape)
            
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_info['table_data'] = self.extract_table_data(shape)
            
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                shape_info['chart_data'] = self.extract_chart_info(shape)
            
            shapes_data.append(shape_info)
        
        return shapes_data
    
    def extract_text_from_shape(self, shape) -> Dict[str, Any]:
        text_data = {
            'full_text': shape.text,
            'paragraphs': []
        }
        
        for paragraph in shape.text_frame.paragraphs:
            para_data = {
                'text': paragraph.text,
                'level': paragraph.level,
                'runs': []
            }
            
            for run in paragraph.runs:
                run_data = {
                    'text': run.text,
                    'font_name': run.font.name,
                    'font_size': run.font.size.pt if run.font.size else None,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'underline': run.font.underline,
                    'color': self.get_color_rgb(run.font.color) if run.font.color else None
                }
                para_data['runs'].append(run_data)
            
            text_data['paragraphs'].append(para_data)
        
        return text_data
    
    def extract_image_info(self, shape) -> Dict[str, Any]:
        try:
            image = shape.image
            return {
                'content_type': image.content_type,
                'filename': image.filename if hasattr(image, 'filename') else 'unknown',
                'width': shape.width,
                'height': shape.height
            }
        except Exception as e:
            logger.warning(f"Could not extract image info: {str(e)}")
            return {}
    
    def extract_table_data(self, shape) -> Dict[str, Any]:
        try:
            table = shape.table
            return {
                'rows': len(table.rows),
                'columns': len(table.columns),
                'cells': [[cell.text for cell in row.cells] for row in table.rows]
            }
        except Exception as e:
            logger.warning(f"Could not extract table data: {str(e)}")
            return {}
    
    def extract_chart_info(self, shape) -> Dict[str, Any]:
        try:
            return {
                'has_chart': True,
                'chart_type': str(shape.chart.chart_type) if hasattr(shape, 'chart') else 'unknown'
            }
        except Exception as e:
            logger.warning(f"Could not extract chart info: {str(e)}")
            return {}
    
    def extract_notes(self, slide) -> str:
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                return notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
            return ""
        except Exception as e:
            logger.warning(f"Could not extract notes: {str(e)}")
            return ""
    
    def extract_branding_info(self) -> Dict[str, Any]:
        branding = {
            'fonts': set(),
            'colors': set(),
            'theme': {}
        }
        
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                branding['fonts'].add(run.font.name)
                            if run.font.color:
                                color = self.get_color_rgb(run.font.color)
                                if color:
                                    branding['colors'].add(color)
        
        branding['fonts'] = list(branding['fonts'])
        branding['colors'] = list(branding['colors'])
        
        return branding
    
    def extract_metadata(self) -> Dict[str, Any]:
        core_props = self.presentation.core_properties
        return {
            'title': core_props.title,
            'author': core_props.author,
            'subject': core_props.subject,
            'slide_count': len(self.presentation.slides),
            'slide_width': self.presentation.slide_width,
            'slide_height': self.presentation.slide_height
        }
    
    def get_color_rgb(self, color) -> str:
        try:
            if hasattr(color, 'rgb'):
                return f"#{color.rgb:06x}"
            return None
        except Exception:
            return None
